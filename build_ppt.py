"""
Build EmPath v2 PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x2B, 0x55)   # dark navy — title bg
TEAL       = RGBColor(0x00, 0x87, 0x8A)   # accent headers
LIGHT_TEAL = RGBColor(0xD6, 0xF0, 0xF0)   # table header bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GREY_BG    = RGBColor(0xF5, 0xF7, 0xFA)
ORANGE     = RGBColor(0xE8, 0x6B, 0x1A)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
RED        = RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

IMG_DIR = "/Users/komalabelursrinivas/Desktop/Capstone/EmPath_v2/Results/error_analysis_v2"


# ── Helper utilities ───────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             bold=False, italic=False, size=18,
             color=BLACK, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.size      = Pt(size)
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb


def navy_header(slide, title, subtitle=None):
    """Full-width navy bar at top with white title."""
    bar_h = Inches(1.35)
    add_rect(slide, 0, 0, W, bar_h, fill=NAVY)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), W - Inches(0.8), Inches(0.85),
             bold=True, size=28, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.88), W - Inches(0.8), Inches(0.42),
                 size=15, color=RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.LEFT)
    # thin teal accent line under bar
    add_rect(slide, 0, bar_h, W, Inches(0.06), fill=TEAL)
    return bar_h + Inches(0.06)


def bullet_block(slide, items, l, t, w, h,
                 size=16, color=BLACK, bullet="▸ ", line_spacing=1.0):
    """Add a bulleted list as a textbox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after  = Pt(0)
        run = p.add_run()
        run.text = bullet + item
        run.font.size      = Pt(size)
        run.font.color.rgb = color
        run.font.name      = "Calibri"
    return txb


def simple_table(slide, headers, rows,
                 l, t, w,
                 header_fill=LIGHT_TEAL, header_color=NAVY,
                 row_fills=None, font_size=13):
    """Draw a simple table using rectangles + text."""
    n_cols = len(headers)
    col_w  = [w // n_cols] * n_cols
    row_h  = Inches(0.38)
    y      = t

    # Header row
    x = l
    for i, hdr in enumerate(headers):
        add_rect(slide, x, y, col_w[i], row_h, fill=header_fill,
                 line=TEAL, line_w=Pt(0.75))
        add_text(slide, hdr, x + Inches(0.05), y + Inches(0.04),
                 col_w[i] - Inches(0.1), row_h - Inches(0.06),
                 bold=True, size=font_size, color=header_color, align=PP_ALIGN.CENTER)
        x += col_w[i]
    y += row_h

    alt_fills = [WHITE, GREY_BG]
    for ri, row in enumerate(rows):
        fill = (row_fills[ri] if row_fills else alt_fills[ri % 2])
        x = l
        for ci, cell in enumerate(row):
            add_rect(slide, x, y, col_w[ci], row_h, fill=fill,
                     line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
            add_text(slide, str(cell),
                     x + Inches(0.05), y + Inches(0.04),
                     col_w[ci] - Inches(0.1), row_h - Inches(0.06),
                     size=font_size - 1, color=BLACK, align=PP_ALIGN.CENTER)
            x += col_w[ci]
        y += row_h

    return y   # bottom y of table


def section_label(slide, text, l, t, w):
    """Teal left-accent section label."""
    add_rect(slide, l, t, Inches(0.08), Inches(0.32), fill=TEAL)
    add_text(slide, text, l + Inches(0.15), t, w, Inches(0.32),
             bold=True, size=14, color=TEAL)


def add_image(slide, path, l, t, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, l, t, w, h)
        else:
            slide.shapes.add_picture(path, l, t, w)


def footer(slide, slide_num, total=20):
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.28), fill=NAVY)
    add_text(slide, f"EmPath v2  |  Komala Belur Srinivas  |  Hofstra University",
             Inches(0.3), H - Inches(0.27), Inches(8), Inches(0.26),
             size=9, color=RGBColor(0xAA, 0xBB, 0xCC))
    add_text(slide, f"{slide_num} / {total}",
             W - Inches(1.2), H - Inches(0.27), Inches(1), Inches(0.26),
             size=9, color=WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(2.8), W, Inches(2.1), fill=TEAL)

add_text(sl, "EmPath v2",
         Inches(0.8), Inches(1.0), W - Inches(1.6), Inches(1.0),
         bold=True, size=52, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Automated Multimodal Pain Intensity Detection",
         Inches(0.8), Inches(2.9), W - Inches(1.6), Inches(0.7),
         bold=True, size=26, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Using Physiological Biosignals & Facial Landmark Geometry",
         Inches(0.8), Inches(3.45), W - Inches(1.6), Inches(0.55),
         size=20, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Komala Belur Srinivas   |   M.S. Computer Science   |   Hofstra University",
         Inches(0.8), Inches(5.2), W - Inches(1.6), Inches(0.42),
         size=16, color=RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
add_text(sl, "BioVid Heat Pain Database  •  LOSO-67 Reactive Subjects  •  PA2 vs PA3",
         Inches(0.8), Inches(5.65), W - Inches(1.6), Inches(0.4),
         size=14, color=RGBColor(0x88, 0xAA, 0xBB), align=PP_ALIGN.CENTER)
footer(sl, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "The Clinical Problem",
                 "Who speaks for patients who cannot speak for themselves?")

add_text(sl, "Pain is subjective. Standard tools require a verbal response — "
             "but millions of patients cannot provide one.",
         Inches(0.4), y0 + Inches(0.12), W - Inches(0.8), Inches(0.5),
         size=15, color=BLACK)

# 4 patient boxes
box_w = Inches(2.9)
box_h = Inches(1.55)
cases = [
    ("🏥  ICU Patients",      "Under sedation or mechanical\nventilation — cannot speak"),
    ("🧠  Dementia Patients",  "Lost verbal communication\ndue to cognitive decline"),
    ("👶  Neonates",           "Cannot self-report pain\nat any developmental stage"),
    ("🔬  Post-Surgical",      "Unconscious or heavily\nmedicated after procedures"),
]
for i, (title, body) in enumerate(cases):
    x = Inches(0.35) + i * (box_w + Inches(0.18))
    add_rect(sl, x, y0 + Inches(0.72), box_w, box_h, fill=WHITE,
             line=TEAL, line_w=Pt(1.5))
    add_rect(sl, x, y0 + Inches(0.72), box_w, Inches(0.42), fill=TEAL)
    add_text(sl, title, x + Inches(0.1), y0 + Inches(0.75),
             box_w - Inches(0.15), Inches(0.38),
             bold=True, size=14, color=WHITE)
    add_text(sl, body, x + Inches(0.1), y0 + Inches(1.18),
             box_w - Inches(0.15), Inches(1.0),
             size=13, color=BLACK)

# Problem statement
add_rect(sl, Inches(0.35), y0 + Inches(2.52), W - Inches(0.7), Inches(1.28),
         fill=NAVY)
add_text(sl, "Current tools (FACS, CPOT, BPS) require trained nurses making manual, "
             "inconsistent bedside observations.\n"
             "Goal: A model that reads biosignals + facial expressions to classify pain intensity "
             "automatically — without the patient saying a word.",
         Inches(0.55), y0 + Inches(2.6), W - Inches(1.1), Inches(1.12),
         size=14, color=WHITE)

section_label(sl, "CLINICAL MOTIVATION", Inches(0.35), y0 + Inches(3.9), Inches(5))
bullet_block(sl,
    ["HIPAA-compliant, non-invasive, real-time capable",
     "Deployable at ICU bedside with existing monitoring hardware",
     "Explainable via SHAP — clinicians can see WHY each prediction was made"],
    Inches(0.5), y0 + Inches(4.25), W - Inches(0.9), Inches(1.1), size=14)
footer(sl, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATASET
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "The Dataset — BioVid Heat Pain Database",
                 "Controlled laboratory pain induction with synchronized biosignals and video")

# Left column — dataset overview
add_rect(sl, Inches(0.35), y0 + Inches(0.2), Inches(5.8), Inches(5.3), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "DATASET OVERVIEW", Inches(0.5), y0 + Inches(0.35), Inches(4))
bullet_block(sl, [
    "87 subjects, heat applied to forearm at controlled temperatures",
    "3 biosignal channels: GSR, ECG, EMG (trapezius + zygomatic)",
    "Video: frontal face camera at 25 fps",
    "Each sample: 5.5-second window → 24 video frames",
    "20 non-reactive subjects excluded (flat biosignals)",
    "67 reactive subjects used for all experiments",
    "~2,680 total samples (PA2 + PA3 balanced)",
], Inches(0.5), y0 + Inches(0.72), Inches(5.5), Inches(2.3), size=14)

# Pain level table
section_label(sl, "PAIN LEVELS", Inches(0.5), y0 + Inches(3.15), Inches(3))
simple_table(sl,
    ["Class", "Temp", "Description"],
    [["BL",  "32°C", "Baseline — no pain"],
     ["PA1", "38°C", "Mild pain"],
     ["PA2", "43°C", "Moderate pain ← used"],
     ["PA3", "45°C", "Intense pain  ← used"],
     ["PA4", "48°C", "Maximum pain"]],
    Inches(0.5), y0 + Inches(3.55), Inches(5.5),
    font_size=13,
    row_fills=[WHITE, GREY_BG, RGBColor(0xFF,0xF3,0xCD),
               RGBColor(0xFF,0xE0,0xCC), GREY_BG])

# Right column — why PA2 vs PA3
add_rect(sl, Inches(6.4), y0 + Inches(0.2), Inches(6.55), Inches(5.3), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "WHY PA2 vs PA3?", Inches(6.55), y0 + Inches(0.35), Inches(5))
add_text(sl, "Only ~1°C thermal separation.\nThe hardest adjacent pair in BioVid literature.",
         Inches(6.55), y0 + Inches(0.72), Inches(6.2), Inches(0.65),
         size=15, bold=True, color=RED)
bullet_block(sl, [
    "PA1 vs PA4 would be trivial — large physiological difference",
    "PA2 vs PA3 is the clinically relevant boundary — distinguishing moderate from intense pain",
    "Published papers confirm this is the most challenging pair",
    "Models that work here generalize to real clinical ambiguity",
], Inches(6.55), y0 + Inches(1.45), Inches(6.2), Inches(2.0), size=14)

section_label(sl, "BIOSIGNAL CHANNELS", Inches(6.55), y0 + Inches(3.55), Inches(5))
simple_table(sl,
    ["Channel", "Measures"],
    [["GSR", "Skin conductance — sympathetic NS response"],
     ["ECG", "Cardiac activity — heart rate + rhythm"],
     ["EMG Trap", "Trapezius — shoulder/neck muscle tension"],
     ["EMG Zyg", "Zygomaticus — cheek smile/pain muscle"]],
    Inches(6.55), y0 + Inches(3.95), Inches(6.3), font_size=12)
footer(sl, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHY THIS IS HARD + KEY DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Why This Is Hard & Key Design Decisions")

# Three difficulty boxes
diff = [
    ("~1°C Thermal Difference",
     "PA2 and PA3 produce nearly identical physiological responses. The model must detect subtle escalation, not obvious pain vs no-pain."),
    ("Massive Inter-Subject Variability",
     "Person A resting GSR = 10 µS. Person B = 2 µS. Absolute signal values cannot be compared across people — each person is their own baseline."),
    ("Tiny Dataset",
     "67 subjects × 40 samples = ~2,680 total. A single ImageNet mini-batch has more data. Deep models simply overfit."),
]
bw = Inches(4.0)
for i, (t, b) in enumerate(diff):
    x = Inches(0.35) + i * (bw + Inches(0.22))
    add_rect(sl, x, y0 + Inches(0.2), bw, Inches(2.05), fill=RED)
    add_text(sl, t, x + Inches(0.15), y0 + Inches(0.25),
             bw - Inches(0.25), Inches(0.55),
             bold=True, size=15, color=WHITE)
    add_text(sl, b, x + Inches(0.15), y0 + Inches(0.82),
             bw - Inches(0.25), Inches(1.3),
             size=13, color=WHITE)

# Design decisions
section_label(sl, "DESIGN DECISION 1 — Exclude 20 Non-Reactive Subjects",
              Inches(0.35), y0 + Inches(2.5), Inches(9))
add_rect(sl, Inches(0.35), y0 + Inches(2.88), W - Inches(0.7), Inches(1.05), fill=WHITE,
         line=TEAL, line_w=Pt(1))
add_text(sl, "20 of 87 subjects show zero physiological response to heat stimulation — "
             "flat GSR, ECG, EMG across all pain levels. Including them means training on pure noise. "
             "Official BioVid evaluation guidelines recommend exclusion.\n"
             "→  67 reactive subjects used for all experiments.",
         Inches(0.55), y0 + Inches(2.95), W - Inches(1.1), Inches(0.92),
         size=14, color=BLACK)

section_label(sl, "DESIGN DECISION 2 — Leave-One-Subject-Out (LOSO) Cross-Validation",
              Inches(0.35), y0 + Inches(4.08), Inches(9))
add_rect(sl, Inches(0.35), y0 + Inches(4.46), W - Inches(0.7), Inches(1.12), fill=WHITE,
         line=TEAL, line_w=Pt(1))
add_text(sl, "67 folds. Each fold: train on 66 subjects → test on 1 held-out subject never seen during training.\n"
             "This tests true cross-subject generalization. Random 80/20 splits leak subject identity — "
             "the model memorizes individual physiology and inflates results by 5–10%.\n"
             "LOSO is the gold standard for clinical pain detection benchmarks.",
         Inches(0.55), y0 + Inches(4.53), W - Inches(1.1), Inches(1.0),
         size=14, color=BLACK)
footer(sl, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Feature Extraction — 35 Biosignal + 22 Landmark Features")

# Biosignal left
add_rect(sl, Inches(0.35), y0 + Inches(0.15), Inches(6.1), Inches(5.45), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "BIOSIGNAL FEATURES (35 total)", Inches(0.5), y0 + Inches(0.3), Inches(5.5))
simple_table(sl,
    ["Channel", "Features"],
    [["GSR",         "mean, std, slope, max, min, range, entropy, sim_corr, sim_mi, gsr×ecg"],
     ["ECG",         "mean, std, max, entropy, sim_corr, sim_mi"],
     ["EMG Trap",    "mean, std, asymmetry (left vs right)"],
     ["HRV (ECG)",   "MeanNN, SDNN, RMSSD, pNN50"]],
    Inches(0.5), y0 + Inches(0.68), Inches(5.85), font_size=12)

add_text(sl, "Tool: neurokit2 (HRV) + numpy (all others)",
         Inches(0.5), y0 + Inches(2.42), Inches(5.8), Inches(0.32),
         size=12, italic=True, color=TEAL)

section_label(sl, "PERSON-SPECIFIC NORMALIZATION", Inches(0.5), y0 + Inches(2.85), Inches(5.5))
add_rect(sl, Inches(0.5), y0 + Inches(3.22), Inches(5.75), Inches(1.5), fill=LIGHT_TEAL)
add_text(sl, "Train subjects: normalize each person using their own mean + std\n"
             "Test subject: normalize using test subject's own statistics\n\n"
             "Person A GSR baseline=10 → PA3=14 (+40%)\n"
             "Person B GSR baseline=2  → PA3=2.8 (+40%)   ← same relative change\n"
             "Without normalization: model sees 14 vs 2.8 — completely different signals",
         Inches(0.65), y0 + Inches(3.28), Inches(5.5), Inches(1.38),
         size=12, color=NAVY)
add_text(sl, "Impact: +3.6% accuracy (59.5% → 63.1%)",
         Inches(0.5), y0 + Inches(4.8), Inches(5.8), Inches(0.32),
         bold=True, size=13, color=GREEN)

# Landmark right
add_rect(sl, Inches(6.65), y0 + Inches(0.15), Inches(6.3), Inches(5.45), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "LANDMARK FEATURES (22 total)", Inches(6.8), y0 + Inches(0.3), Inches(5.5))
add_text(sl, "MediaPipe FaceMesh → 468 (x,y,z) landmarks per frame × 24 frames per sample",
         Inches(6.8), y0 + Inches(0.68), Inches(6.0), Inches(0.45),
         size=13, color=BLACK)
simple_table(sl,
    ["Feature Group", "Features (mean + std)"],
    [["Eye Openness",  "left, right, avg eye opening"],
     ["Brow Position", "left, right, avg brow-to-eye distance"],
     ["Brow Furrow",   "inter-brow distance"],
     ["Mouth",         "height, width, aspect ratio"],
     ["Nose",          "nose width"]],
    Inches(6.8), y0 + Inches(1.18), Inches(6.0), font_size=12)

add_text(sl, "→ mean + std across 24 frames = 22 total features",
         Inches(6.8), y0 + Inches(3.22), Inches(6.0), Inches(0.35),
         bold=True, size=13, color=TEAL)

section_label(sl, "KEY INSIGHT FROM SHAP", Inches(6.8), y0 + Inches(3.65), Inches(5.5))
add_rect(sl, Inches(6.8), y0 + Inches(4.02), Inches(6.0), Inches(1.2), fill=NAVY)
add_text(sl, "All top landmark features are _std (variability), not _mean (position).\n"
             "Pain expression is about facial dynamics — how much the mouth and brows\n"
             "move across the 5.5s window — not where they sit statically.",
         Inches(6.95), y0 + Inches(4.1), Inches(5.7), Inches(1.05),
         size=13, color=WHITE)
footer(sl, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BASELINE MODEL
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Baseline — EmPath Stacked Fusion Architecture",
                 "RF(biosignal) + RF(landmark) → LogisticRegression meta-learner")

# Architecture diagram using shapes
mid = Inches(6.66)
bx = Inches(0.5); by = y0 + Inches(0.4); bw2 = Inches(3.6); bh = Inches(0.65)

add_rect(sl, bx, by, bw2, bh, fill=RGBColor(0x00,0x60,0x80)); add_text(sl, "Biosignal Features (35)", bx+Inches(0.1), by+Inches(0.1), bw2-Inches(0.15), bh-Inches(0.15), bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, W-bx-bw2, by, bw2, bh, fill=RGBColor(0x1B,0x6B,0x2E)); add_text(sl, "Landmark Features (22)", W-bx-bw2+Inches(0.1), by+Inches(0.1), bw2-Inches(0.15), bh-Inches(0.15), bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)

arrow_y1 = by + bh + Inches(0.08)
add_text(sl, "↓", bx+Inches(1.4), arrow_y1, Inches(0.8), Inches(0.4), size=22, color=TEAL, align=PP_ALIGN.CENTER)
add_text(sl, "↓", W-bx-bw2+Inches(1.4), arrow_y1, Inches(0.8), Inches(0.4), size=22, color=TEAL, align=PP_ALIGN.CENTER)

rf_y = arrow_y1 + Inches(0.42)
add_rect(sl, bx, rf_y, bw2, bh, fill=RGBColor(0x00,0x60,0x80)); add_text(sl, "RF Biosignal\n(300 trees, depth 4)", bx+Inches(0.1), rf_y+Inches(0.04), bw2-Inches(0.15), bh-Inches(0.08), bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, W-bx-bw2, rf_y, bw2, bh, fill=RGBColor(0x1B,0x6B,0x2E)); add_text(sl, "RF Landmark\n(300 trees, depth 4)", W-bx-bw2+Inches(0.1), rf_y+Inches(0.04), bw2-Inches(0.15), bh-Inches(0.08), bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)

arrow_y2 = rf_y + bh + Inches(0.08)
add_text(sl, "↓  P(PA2), P(PA3)", bx+Inches(0.5), arrow_y2, Inches(2.8), Inches(0.4), size=13, color=TEAL)
add_text(sl, "↓  P(PA2), P(PA3)", W-bx-bw2+Inches(0.5), arrow_y2, Inches(2.8), Inches(0.4), size=13, color=GREEN)

concat_y = arrow_y2 + Inches(0.42)
add_rect(sl, mid - Inches(2.2), concat_y, Inches(4.4), Inches(0.48), fill=TEAL)
add_text(sl, "Concatenate → [P_bio_PA2, P_bio_PA3, P_lm_PA2, P_lm_PA3]",
         mid - Inches(2.1), concat_y + Inches(0.06), Inches(4.2), Inches(0.38),
         bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "↓", mid - Inches(0.2), concat_y + Inches(0.5), Inches(0.4), Inches(0.38), size=22, color=TEAL, align=PP_ALIGN.CENTER)

lr_y = concat_y + Inches(0.9)
add_rect(sl, mid - Inches(2.2), lr_y, Inches(4.4), bh, fill=NAVY)
add_text(sl, "Logistic Regression (meta-learner)", mid - Inches(2.1), lr_y + Inches(0.1), Inches(4.2), bh - Inches(0.15), bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "↓", mid - Inches(0.2), lr_y + bh + Inches(0.04), Inches(0.4), Inches(0.38), size=22, color=ORANGE, align=PP_ALIGN.CENTER)

pred_y = lr_y + bh + Inches(0.45)
add_rect(sl, mid - Inches(2.0), pred_y, Inches(4.0), bh, fill=ORANGE)
add_text(sl, "Final Prediction: PA2 or PA3", mid - Inches(1.9), pred_y + Inches(0.1), Inches(3.8), bh - Inches(0.15), bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)

# Right side — why choices
rx = W - Inches(4.5); ry = y0 + Inches(0.4)
add_rect(sl, rx, ry, Inches(4.1), Inches(5.5), fill=WHITE, line=TEAL, line_w=Pt(1))
section_label(sl, "WHY THESE CHOICES?", rx + Inches(0.15), ry + Inches(0.1), Inches(3.8))
bullet_block(sl, [
    "RF beats deep learning at N=67 — shallow depth prevents overfitting",
    "min_samples_split=10 forces generalizable splits",
    "Stacked (65.3%) > Early fusion (64.6%) — calibrated probs more informative than raw features on different scales",
    "LogReg meta-learner is regularized, fast, interpretable",
    "Person-norm applied inside each fold — zero leakage",
], rx + Inches(0.15), ry + Inches(0.52), Inches(3.8), Inches(3.0), size=13)

add_rect(sl, rx, pred_y - Inches(0.1), Inches(4.1), Inches(1.1), fill=NAVY)
add_text(sl, "RESULT\n65.3% ± 14.1%  LOSO-67",
         rx + Inches(0.15), pred_y - Inches(0.05), Inches(3.85), Inches(1.0),
         bold=True, size=18, color=ORANGE, align=PP_ALIGN.CENTER)
footer(sl, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ALL EXPERIMENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "15+ Architectures Systematically Evaluated")

exp_rows = [
    ["Vision MobileNetV2",          "Random split", "47.2%"],
    ["Biosignal SVM",               "Random split", "48.8%"],
    ["Biosignal MLP",               "Random split", "51.2%"],
    ["Landmark RF (flat)",          "Random split", "51.6%"],
    ["Biosignal XGBoost",           "Random split", "54.1%"],
    ["BIOT Foundation Model",       "LOSO-67",      "54.4%"],
    ["Biosignal TCN",               "Random split", "55.9%"],
    ["PainFormer",                  "LOSO-67",      "53.1%"],
    ["Tiny-BioMoE",                 "LOSO-67",      "56.7%"],
    ["Biosignal RF (67 subjects)",  "Random split", "59.5%"],
    ["Biosignal RF + person-norm",  "LOSO-67",      "63.1%"],
    ["Landmark RF (flat)",          "LOSO-67",      "61.4%"],
    ["Early fusion (concat)",       "LOSO-67",      "64.6%"],
    ["Subject adaptation RF",       "LOSO-67",      "65.1%"],
    ["CORAL ordinal MLP",           "LOSO-67",      "65.3%"],
    ["EmPath Stacked Fusion ★",     "LOSO-67",      "65.3% ± 14.1%"],
]
row_fills_exp = []
for r in exp_rows:
    if "★" in r[0]:
        row_fills_exp.append(ORANGE)
    elif r[2].startswith("6") and "LOSO" in r[1]:
        row_fills_exp.append(LIGHT_TEAL)
    elif "Random" in r[1]:
        row_fills_exp.append(RGBColor(0xF5,0xF5,0xF5))
    else:
        row_fills_exp.append(WHITE)

simple_table(sl,
    ["Model", "Protocol", "Accuracy"],
    exp_rows,
    Inches(0.35), y0 + Inches(0.2), W - Inches(0.7),
    font_size=12, row_fills=row_fills_exp)
footer(sl, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GNN
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Novel Contribution 1 — Graph Attention Network on Facial Landmarks",
                 "Modeling anatomical co-activation in facial pain expression")

# Left: architecture description
add_rect(sl, Inches(0.35), y0 + Inches(0.2), Inches(6.0), Inches(5.4), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "HYPOTHESIS", Inches(0.5), y0 + Inches(0.35), Inches(5))
add_text(sl, "22 flat statistics discard spatial relationships. A GNN learns relational "
             "co-activation: 'brow lowers simultaneously with lip tightening'.",
         Inches(0.5), y0 + Inches(0.7), Inches(5.7), Inches(0.65), size=13, color=BLACK)

section_label(sl, "ARCHITECTURE", Inches(0.5), y0 + Inches(1.45), Inches(5))
bullet_block(sl, [
    "Nodes: 468 MediaPipe landmarks, features = (x,y) coordinates",
    "Edges: Delaunay triangulation → 1,365 edges per frame",
    "Per-frame normalization: nose-centered, IOD-scaled",
    "2-layer GAT (Graph Attention), 4 attention heads",
    "Global mean pool → 32-dim embedding per frame",
    "Average 24 frames → sample embedding",
    "Stack with RF_biosignal probs → LogReg final",
], Inches(0.5), y0 + Inches(1.82), Inches(5.7), Inches(2.3), size=13)

section_label(sl, "IMPLEMENTATION CHALLENGE", Inches(0.5), y0 + Inches(4.2), Inches(5))
add_rect(sl, Inches(0.5), y0 + Inches(4.55), Inches(5.7), Inches(0.9), fill=LIGHT_TEAL)
add_text(sl, "Batch.from_data_list() — 5 min/fold (too slow for 67 folds).\n"
             "Fix: GPU-cached vectorized batching — precomputed edge offsets on GPU, "
             "eliminated 33MB edge_index transfer per batch → ~90 sec/fold.",
         Inches(0.65), y0 + Inches(4.62), Inches(5.4), Inches(0.8), size=12, color=NAVY)

# Right: results + why failed
add_rect(sl, Inches(6.6), y0 + Inches(0.2), Inches(6.35), Inches(2.4), fill=WHITE,
         line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "RESULTS", Inches(6.75), y0 + Inches(0.35), Inches(5))
simple_table(sl,
    ["Model", "Accuracy", "Std"],
    [["GNN landmarks only",      "51.7%", "±9.8%"],
     ["GNN + RF biosignal",      "63.1%", "±11.9%"],
     ["Baseline stacked fusion", "65.3%", "±14.1%"]],
    Inches(6.75), y0 + Inches(0.72), Inches(6.0), font_size=13,
    row_fills=[RED, LIGHT_TEAL, ORANGE])

add_rect(sl, Inches(6.6), y0 + Inches(2.75), Inches(6.35), Inches(2.85), fill=NAVY)
section_label(sl, "WHY IT FAILED TO BEAT BASELINE", Inches(6.75), y0 + Inches(2.9), Inches(5.5))
add_text(sl, "  ", Inches(6.75), y0 + Inches(2.9), Inches(0.1), Inches(0.1), size=1, color=WHITE)
bullet_block(sl, [
    "67 subjects is too few to train a GNN to learn spatial representations from scratch",
    "Graph topology learning requires many examples per edge pattern",
    "RF has built-in structural inductive bias; GAT must learn it from data",
    "This is a sample size bottleneck, not an architectural failure",
    "With N=500+ subjects, GNN would likely outperform flat features",
], Inches(6.75), y0 + Inches(3.28), Inches(6.0), Inches(2.1), size=13, color=WHITE)

add_rect(sl, Inches(6.6), y0 + Inches(5.7), Inches(6.35), Inches(0.65), fill=TEAL)
add_text(sl, "No published BioVid paper uses GNNs on MediaPipe FaceMesh — genuinely novel methodology",
         Inches(6.75), y0 + Inches(5.77), Inches(6.05), Inches(0.52),
         size=13, bold=True, color=WHITE)
footer(sl, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DANN
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Novel Contribution 2 — DANN Domain Adaptation",
                 "Adversarial training for subject-invariant pain features")

section_label(sl, "HYPOTHESIS", Inches(0.35), y0 + Inches(0.2), Inches(6))
add_text(sl, "±14.1% std means the model is inconsistent across subjects. "
             "Force the encoder to learn features that are informative for pain "
             "but useless for identifying the subject.",
         Inches(0.35), y0 + Inches(0.55), Inches(7.5), Inches(0.55), size=14, color=BLACK)

# Architecture diagram
ax, ay = Inches(0.5), y0 + Inches(1.2)
add_rect(sl, ax, ay, Inches(3.0), Inches(0.5), fill=RGBColor(0x00,0x60,0x80))
add_text(sl, "Biosignal Features (35)", ax+Inches(0.05), ay+Inches(0.07), Inches(2.9), Inches(0.38), bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "↓", ax+Inches(1.2), ay+Inches(0.52), Inches(0.6), Inches(0.38), size=20, color=TEAL, align=PP_ALIGN.CENTER)
ay2 = ay + Inches(0.92)
add_rect(sl, ax, ay2, Inches(3.0), Inches(0.55), fill=TEAL)
add_text(sl, "Shared FC Encoder\n35 → 64 → 32", ax+Inches(0.05), ay2+Inches(0.04), Inches(2.9), Inches(0.48), bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
# Fork
add_text(sl, "↓", ax+Inches(0.3), ay2+Inches(0.57), Inches(0.6), Inches(0.38), size=20, color=GREEN, align=PP_ALIGN.CENTER)
add_text(sl, "↓", ax+Inches(1.8), ay2+Inches(0.57), Inches(0.6), Inches(0.38), size=20, color=RED, align=PP_ALIGN.CENTER)
ay3 = ay2 + Inches(0.97)
add_rect(sl, ax, ay3, Inches(1.3), Inches(0.72), fill=GREEN)
add_text(sl, "Pain\nClassifier\nFC(32→2)", ax+Inches(0.05), ay3+Inches(0.04), Inches(1.2), Inches(0.65), bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, ax+Inches(1.6), ay3, Inches(1.4), Inches(0.72), fill=RED)
add_text(sl, "Gradient\nReversal\nLayer (GRL)", ax+Inches(1.65), ay3+Inches(0.04), Inches(1.3), Inches(0.65), bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "↓", ax+Inches(2.1), ay3+Inches(0.74), Inches(0.6), Inches(0.35), size=18, color=RED, align=PP_ALIGN.CENTER)
ay4 = ay3 + Inches(1.1)
add_rect(sl, ax+Inches(1.5), ay4, Inches(1.5), Inches(0.72), fill=RED)
add_text(sl, "Subject\nClassifier\nFC(32→67)", ax+Inches(1.55), ay4+Inches(0.04), Inches(1.4), Inches(0.65), bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# GRL explanation
add_rect(sl, Inches(3.85), y0 + Inches(1.2), Inches(4.2), Inches(3.3), fill=WHITE, line=TEAL, line_w=Pt(1))
section_label(sl, "GRADIENT REVERSAL LAYER", Inches(4.0), y0 + Inches(1.35), Inches(3.8))
add_text(sl, "Forward pass: identity function — passes features through unchanged\n\n"
             "Backward pass: multiply gradient by −α\n\n"
             "Effect: encoder learns features that maximize pain classification accuracy "
             "AND minimize subject identification accuracy simultaneously.\n\n"
             "Only 5 lines of PyTorch (torch.autograd.Function).\n\n"
             "α ramps 0→1 over training — starts with pure pain learning, "
             "gradually adds adversarial domain pressure.",
         Inches(4.0), y0 + Inches(1.72), Inches(3.9), Inches(2.7), size=12, color=BLACK)

# Results
add_rect(sl, Inches(8.3), y0 + Inches(1.2), Inches(4.65), Inches(2.1), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "RESULTS", Inches(8.45), y0 + Inches(1.35), Inches(4))
simple_table(sl,
    ["Model", "Accuracy", "Std"],
    [["DANN biosignal only",      "61.6%", "±10.3%"],
     ["DANN + RF landmarks",      "64.7%", "±11.8%"],
     ["Baseline stacked fusion",  "65.3%", "±14.1%"]],
    Inches(8.45), y0 + Inches(1.72), Inches(4.35), font_size=12,
    row_fills=[LIGHT_TEAL, LIGHT_TEAL, ORANGE])

add_rect(sl, Inches(8.3), y0 + Inches(3.45), Inches(4.65), Inches(1.9), fill=NAVY)
add_text(sl, "KEY FINDING\n\n"
             "Std dropped from ±14.1% → ±10.3% (−3.8%).\n"
             "Adversarial training makes predictions more CONSISTENT across subjects "
             "but does not improve mean accuracy — the bottleneck is sample size, not domain shift.",
         Inches(8.45), y0 + Inches(3.55), Inches(4.35), Inches(1.75),
         size=12, color=WHITE)

# Bottom full-width insight
add_rect(sl, Inches(0.35), y0 + Inches(5.42), W - Inches(0.7), Inches(0.72), fill=TEAL)
add_text(sl, "DANN confirmed: subject identity IS encoded in biosignals. "
             "Removing it stabilizes predictions but insufficient signal remains to surpass the RF ceiling at N=67.",
         Inches(0.55), y0 + Inches(5.52), W - Inches(1.1), Inches(0.58),
         size=13, bold=True, color=WHITE)
footer(sl, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CrossMod + Velocity (combined)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Novel Contributions 3 & 4 — CrossMod Attention + Landmark Velocity")

# CrossMod left
add_rect(sl, Inches(0.35), y0 + Inches(0.2), Inches(6.2), Inches(5.45), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
add_rect(sl, Inches(0.35), y0 + Inches(0.2), Inches(6.2), Inches(0.38), fill=TEAL)
add_text(sl, "CrossMod — Bidirectional Cross-Modal Attention", Inches(0.5), y0 + Inches(0.25), Inches(6.0), Inches(0.3), bold=True, size=14, color=WHITE)
add_text(sl, "Previous attention fusion (61.1%) used self-attention within each modality. "
             "CrossMod uses biosignals attending to landmarks and vice versa.",
         Inches(0.5), y0 + Inches(0.72), Inches(5.9), Inches(0.55), size=13, color=BLACK)
bullet_block(sl, [
    "bio_emb queries lm_emb → 'which landmark patterns confirm the biosignal?'",
    "lm_emb queries bio_emb → 'which physiological patterns confirm the face?'",
    "Residual + LayerNorm on both attended outputs",
    "Concatenate (128-dim) → FC → pain classifier",
    "4 attention heads, 100 epochs, cosine LR schedule",
], Inches(0.5), y0 + Inches(1.35), Inches(5.9), Inches(2.0), size=13)

section_label(sl, "RESULTS", Inches(0.5), y0 + Inches(3.45), Inches(2))
simple_table(sl,
    ["Model", "Acc", "Std"],
    [["CrossMod",  "63.1%", "±11.1%"],
     ["Baseline",  "65.3%", "±14.1%"]],
    Inches(0.5), y0 + Inches(3.82), Inches(5.8), font_size=13,
    row_fills=[LIGHT_TEAL, ORANGE])
add_text(sl, "−2.2% accuracy but −3.0% std. Cross-attention improves consistency. "
             "RF calibrated probs beat learned attention at N=67.",
         Inches(0.5), y0 + Inches(4.75), Inches(5.9), Inches(0.65), size=12, italic=True, color=BLACK)

# Velocity right
add_rect(sl, Inches(6.8), y0 + Inches(0.2), Inches(6.15), Inches(5.45), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
add_rect(sl, Inches(6.8), y0 + Inches(0.2), Inches(6.15), Inches(0.38), fill=GREEN)
add_text(sl, "Landmark Velocity — Optical Flow Proxy", Inches(6.95), y0 + Inches(0.25), Inches(5.9), Inches(0.3), bold=True, size=14, color=WHITE)
add_text(sl, "If PA3 produces faster facial movements than PA2, frame-to-frame displacement vectors "
             "should carry discriminative signal that static geometry misses.",
         Inches(6.95), y0 + Inches(0.72), Inches(5.8), Inches(0.55), size=13, color=BLACK)

add_rect(sl, Inches(6.95), y0 + Inches(1.35), Inches(5.8), Inches(1.05), fill=LIGHT_TEAL)
add_text(sl, "velocity[t] = position[t+1] − position[t]  →  (N, 23, 468, 2)\n"
             "magnitude = ||velocity||₂  →  (N, 23, 468)\n"
             "Per-frame stats: mean, std, max across 468 landmarks  →  3×23 = 69 features\n"
             "Global stats: mean, std, max across all frames  →  3 features\n"
             "Total: 72 velocity features per sample. No video reprocessing needed.",
         Inches(6.95), y0 + Inches(1.4), Inches(5.8), Inches(0.95), size=12, color=NAVY)

section_label(sl, "RESULTS", Inches(6.95), y0 + Inches(2.55), Inches(2))
simple_table(sl,
    ["Model", "Acc", "Std"],
    [["Velocity RF only",       "60.0%", "±11.9%"],
     ["Velocity + biosignal",   "64.0%", "±12.7%"],
     ["Static landmark RF",     "61.4%", "±13.1%"],
     ["Baseline stacked",       "65.3%", "±14.1%"]],
    Inches(6.95), y0 + Inches(2.92), Inches(5.8), font_size=12,
    row_fills=[LIGHT_TEAL, LIGHT_TEAL, GREY_BG, ORANGE])

add_rect(sl, Inches(6.8), y0 + Inches(4.8), Inches(6.15), Inches(0.78), fill=NAVY)
add_text(sl, "Velocity features WORSE than static positions.\n"
             "PA3 is a HELD facial expression, not a fast movement. "
             "Pain signal lives in sustained geometry, not rapid motion.",
         Inches(6.95), y0 + Inches(4.88), Inches(5.9), Inches(0.65),
         size=13, color=WHITE, bold=True)
footer(sl, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SHAP Analysis
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Novel Contribution 5 — SHAP / XAI Explainability",
                 "Out-of-fold SHAP TreeExplainer on both RF models — no data leakage")

section_label(sl, "WHY XAI FOR CLINICAL DEPLOYMENT", Inches(0.35), y0 + Inches(0.2), Inches(7))
add_text(sl, "A 65% accurate black box is not deployable in an ICU. "
             "Clinicians need to know WHICH signals drive each prediction before trusting it.",
         Inches(0.35), y0 + Inches(0.55), W - Inches(0.7), Inches(0.42), size=14, color=BLACK)

# Biosignal SHAP table
add_rect(sl, Inches(0.35), y0 + Inches(1.1), Inches(5.9), Inches(4.58), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "TOP BIOSIGNAL FEATURES (by SHAP)", Inches(0.5), y0 + Inches(1.25), Inches(5.5))
simple_table(sl,
    ["Rank", "Feature", "SHAP", "Clinical Meaning"],
    [["1", "gsr_slope",     "0.037", "Rate of GSR rise — sympathetic NS activation speed"],
     ["2", "gsr_std",       "0.028", "GSR variability — autonomic turbulence under pain"],
     ["3", "ecg_max",       "0.020", "Peak cardiac amplitude — heart rate spike"],
     ["4", "gsr_shannon",   "0.014", "GSR entropy — arousal signal complexity"],
     ["5", "ecg_shannon",   "0.014", "Cardiac entropy — cardiac arousal complexity"],
     ["6", "gsr_sim_corr",  "0.014", "GSR similarity to pain template correlation"],
     ["9", "emg_trap_std",  "0.008", "Trapezius variability — shoulder tension"],
     ["12","hrv_meannn",    "0.006", "Mean R-R interval — cardiac autonomic modulation"]],
    Inches(0.5), y0 + Inches(1.62), Inches(5.65), font_size=11)

add_rect(sl, Inches(0.5), y0 + Inches(4.9), Inches(5.65), Inches(0.62), fill=TEAL)
add_text(sl, "GSR dominates — skin conductance is the most sensitive autonomic pain marker. "
             "Rate of rise (slope) matters more than absolute level.",
         Inches(0.65), y0 + Inches(4.97), Inches(5.35), Inches(0.52), size=12, bold=True, color=WHITE)

# Landmark SHAP table
add_rect(sl, Inches(6.5), y0 + Inches(1.1), Inches(6.45), Inches(4.58), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "TOP LANDMARK FEATURES (by SHAP)", Inches(6.65), y0 + Inches(1.25), Inches(5.5))
simple_table(sl,
    ["Rank", "Feature", "SHAP", "Clinical Meaning"],
    [["1", "mouth_height_std",      "0.029", "Lip parting variability across 5.5s window"],
     ["2", "mouth_width_std",       "0.025", "Lateral lip movement variability"],
     ["3", "nose_width_std",        "0.025", "Nostril flaring variability"],
     ["4", "mouth_aspect_ratio_std","0.016", "Overall mouth shape change dynamics"],
     ["5", "left_brow_eye_dist_std","0.013", "Brow raise/lower variability"],
     ["6", "brow_eye_avg_std",      "0.012", "Average brow position variability"],
     ["7", "avg_eye_openness_mean", "0.009", "Mean eye opening — sustained expression"],
     ["8", "brow_furrow_std",       "0.007", "Inter-brow distance variability"]],
    Inches(6.65), y0 + Inches(1.62), Inches(6.15), font_size=11)

add_rect(sl, Inches(6.65), y0 + Inches(4.9), Inches(6.15), Inches(0.62), fill=GREEN)
add_text(sl, "All top features are _std (variability). Pain expression is about facial DYNAMICS, "
             "not static geometry. Validates the stacked fusion design.",
         Inches(6.8), y0 + Inches(4.97), Inches(5.85), Inches(0.52), size=12, bold=True, color=WHITE)
footer(sl, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SHAP Plots
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "SHAP Visualizations — Feature Impact on PA3 Prediction")

img_w = Inches(6.1)
img_h = Inches(5.5)

add_image(sl,
    os.path.join(IMG_DIR, "shap_biosignal_bar.png"),
    Inches(0.2), y0 + Inches(0.15), img_w, img_h)

add_image(sl,
    os.path.join(IMG_DIR, "shap_biosignal_beeswarm.png"),
    Inches(6.85), y0 + Inches(0.15), img_w, img_h)
footer(sl, 12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Landmark SHAP + Per-Subject
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Landmark SHAP & Per-Subject Accuracy Breakdown")

add_image(sl,
    os.path.join(IMG_DIR, "shap_landmark_bar.png"),
    Inches(0.2), y0 + Inches(0.15), Inches(5.8), Inches(3.5))

add_image(sl,
    os.path.join(IMG_DIR, "per_subject_accuracy.png"),
    Inches(0.2), y0 + Inches(3.75), W - Inches(0.35), Inches(3.0))

add_rect(sl, Inches(6.2), y0 + Inches(0.15), Inches(6.75), Inches(3.5), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "PER-SUBJECT BREAKDOWN", Inches(6.35), y0 + Inches(0.3), Inches(5))

stats = [("≥ 80%", "14 subjects", GREEN, "Model generalizes — consistent physiology"),
         ("65–80%", "23 subjects", TEAL, "Above chance, reliable predictions"),
         ("50–65%", "23 subjects", ORANGE, "Near chance, minimal separable signal"),
         ("< 50%", "7 subjects", RED, "Below chance — atypical response patterns")]
sy = y0 + Inches(0.7)
for rng, n, col, note in stats:
    add_rect(sl, Inches(6.35), sy, Inches(1.1), Inches(0.5), fill=col)
    add_text(sl, rng, Inches(6.4), sy + Inches(0.06), Inches(1.0), Inches(0.4), bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, n, Inches(7.55), sy + Inches(0.06), Inches(1.2), Inches(0.4), bold=True, size=14, color=col)
    add_text(sl, note, Inches(8.8), sy + Inches(0.1), Inches(4.0), Inches(0.35), size=12, color=BLACK)
    sy += Inches(0.62)

add_rect(sl, Inches(6.35), sy + Inches(0.15), Inches(6.5), Inches(1.1), fill=NAVY)
add_text(sl, "The bimodal distribution (strong right cluster vs weak left) reveals SUBJECT "
             "HETEROGENEITY — not model failure. Some people's PA2 and PA3 are physiologically "
             "indistinguishable even to human experts. This is a biological constraint.",
         Inches(6.5), sy + Inches(0.22), Inches(6.2), Inches(0.98), size=12, color=WHITE)
footer(sl, 13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Confusion Matrix + Feature Importance
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Confusion Matrix & Combined Feature Importance")

add_image(sl,
    os.path.join(IMG_DIR, "confusion_matrix.png"),
    Inches(0.3), y0 + Inches(0.2), Inches(5.2), Inches(4.3))

add_image(sl,
    os.path.join(IMG_DIR, "feature_importance_combined.png"),
    Inches(5.8), y0 + Inches(0.2), Inches(7.2), Inches(4.3))

# Interpretation row
add_rect(sl, Inches(0.3), y0 + Inches(4.6), Inches(4.9), Inches(1.05), fill=WHITE, line=TEAL, line_w=Pt(1))
add_text(sl, "Balanced errors: 456 PA2→PA3 vs 470 PA3→PA2.\n"
             "Model not biased toward either class — struggling equally in both directions, "
             "as expected with only ~1°C thermal separation.",
         Inches(0.45), y0 + Inches(4.67), Inches(4.65), Inches(0.92), size=13, color=BLACK)

add_rect(sl, Inches(5.5), y0 + Inches(4.6), Inches(7.5), Inches(1.05), fill=WHITE, line=GREEN, line_w=Pt(1))
add_text(sl, "gsr_slope dominates (0.11). Landmark _std features occupy 4 of top 6 positions "
             "(mouth_width_std, mouth_aspect_ratio_std, mouth_height_std). "
             "Both modalities contribute non-redundantly — validates the stacked fusion design.",
         Inches(5.65), y0 + Inches(4.67), Inches(7.2), Inches(0.92), size=13, color=BLACK)
footer(sl, 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Why No Novel Method Beat Baseline
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Why No Novel Method Beat the Baseline — Honest Scientific Finding")

add_rect(sl, Inches(0.35), y0 + Inches(0.2), W - Inches(0.7), Inches(0.65), fill=NAVY)
add_text(sl, "The RF stacked fusion is the empirical ceiling for N=67 subjects. "
             "This is not architectural failure — it is sample size.",
         Inches(0.55), y0 + Inches(0.3), W - Inches(1.1), Inches(0.5),
         bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)

# Why each failed
fails = [
    ("GNN",     "51.7% ± 9.8%",  "Graph topology learning requires many examples per edge pattern. 67 subjects = too few to learn spatial representations from scratch."),
    ("DANN",    "61.6% ± 10.3%", "Adversarial training reduces std by 3.8% but cannot increase mean accuracy — subject heterogeneity is biological, not a learnable domain."),
    ("CrossMod","63.1% ± 11.1%", "Learned cross-modal attention cannot outperform calibrated RF probabilities when training set = 66 subjects per fold."),
    ("Velocity","60.0% ± 11.9%", "PA3 is a HELD expression, not a fast movement. Dynamic motion features carry less signal than sustained geometric variability."),
]
fy = y0 + Inches(1.05)
for name, acc, reason in fails:
    add_rect(sl, Inches(0.35), fy, Inches(1.5), Inches(0.9), fill=RED)
    add_text(sl, name, Inches(0.4), fy + Inches(0.1), Inches(1.4), Inches(0.72), bold=True, size=17, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.88), fy, Inches(1.9), Inches(0.9), fill=LIGHT_TEAL)
    add_text(sl, acc, Inches(1.93), fy + Inches(0.15), Inches(1.8), Inches(0.6), bold=True, size=15, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.88), fy, W - Inches(4.22), Inches(0.9), fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
    add_text(sl, reason, Inches(4.03), fy + Inches(0.1), W - Inches(4.5), Inches(0.75), size=13, color=BLACK)
    fy += Inches(1.0)

# What they DID achieve
add_rect(sl, Inches(0.35), fy + Inches(0.1), W - Inches(0.7), Inches(0.42), fill=TEAL)
add_text(sl, "WHAT THE NOVEL METHODS DID ACHIEVE", Inches(0.55), fy + Inches(0.15), W - Inches(1.0), Inches(0.32), bold=True, size=14, color=WHITE)
fy2 = fy + Inches(0.55)
add_rect(sl, Inches(0.35), fy2, W - Inches(0.7), Inches(0.78), fill=WHITE, line=TEAL, line_w=Pt(1))
bullet_block(sl, [
    "Consistently lower standard deviation (±10–12% vs ±14.1%) — more reliable, consistent predictions across subjects",
    "Confirmed domain adaptation reduces subject-identity leakage  |  Confirmed velocity < static variability (biological insight)",
    "Demonstrated that cross-modal attention matches but cannot surpass RF at this scale — directionally correct, scale-limited",
], Inches(0.55), fy2 + Inches(0.06), W - Inches(1.0), Inches(0.68), size=13)
footer(sl, 15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Complete Results Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Complete Results Summary — All Novel Contributions")

all_rows = [
    ["Biosignal RF (baseline)",         "LOSO-67", "63.1%", "±11.6%", "—"],
    ["Landmark RF flat (baseline)",     "LOSO-67", "61.4%", "±13.1%", "—"],
    ["Stacked Fusion ★ (baseline)",     "LOSO-67", "65.3%", "±14.1%", "Baseline"],
    ["GNN landmarks only",              "LOSO-67", "51.7%", "±9.8%",  "Novel"],
    ["GNN + RF biosignal stacked",      "LOSO-67", "63.1%", "±11.9%", "Novel"],
    ["DANN biosignal",                  "LOSO-67", "61.6%", "±10.3%", "Novel ↓std"],
    ["DANN + RF landmarks stacked",     "LOSO-67", "64.7%", "±11.8%", "Novel ↓std"],
    ["CrossMod cross-attention",        "LOSO-67", "63.1%", "±11.1%", "Novel ↓std"],
    ["Velocity RF only",                "LOSO-67", "60.0%", "±11.9%", "Novel"],
    ["Velocity + biosignal stacked",    "LOSO-67", "64.0%", "±12.7%", "Novel"],
]
fills_all = [GREY_BG, GREY_BG, ORANGE,
             LIGHT_TEAL, LIGHT_TEAL,
             RGBColor(0xE8,0xF5,0xE9), RGBColor(0xE8,0xF5,0xE9),
             RGBColor(0xE3,0xF2,0xFD),
             LIGHT_TEAL, LIGHT_TEAL]

simple_table(sl,
    ["Model", "Protocol", "Accuracy", "Std Dev", "Type"],
    all_rows,
    Inches(0.35), y0 + Inches(0.2), W - Inches(0.7),
    font_size=13, row_fills=fills_all)
footer(sl, 16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Literature Comparison
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Literature Comparison & Honest Gap Analysis")

simple_table(sl,
    ["Method", "Eval Protocol", "Subjects", "Accuracy", "Notes"],
    [["Biosignal SVM",           "Random split",         "50",  "48.8%", "Flat baseline, no cross-subject test"],
     ["Hybrid CNN + features",   "LOSO-67 reactive",     "67",  "59.7%", "Prior work"],
     ["EmPath Stacked Fusion ★", "LOSO-67 reactive",     "67",  "65.3%", "Our best — stricter evaluation"],
     ["CrossMod-Transformer 2025","LOSO-87 all subjects","87",  "87.5%", "Includes 20 non-reactive subjects"]],
    Inches(0.35), y0 + Inches(0.25), W - Inches(0.7),
    font_size=14,
    row_fills=[GREY_BG, GREY_BG, ORANGE, LIGHT_TEAL])

add_rect(sl, Inches(0.35), y0 + Inches(2.48), W - Inches(0.7), Inches(0.45), fill=RED)
add_text(sl, "22% absolute gap to published SOTA — this is real and must be acknowledged",
         Inches(0.55), y0 + Inches(0.55 + 1.98), W - Inches(1.0), Inches(0.35),
         bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)

section_label(sl, "THE HONEST DEFENSE", Inches(0.35), y0 + Inches(3.1), Inches(5))
defenses = [
    "CrossMod-Transformer evaluates on ALL 87 subjects including the 20 non-reactive ones whose biosignals are flat — making the task trivially easier",
    "We use REACTIVE-ONLY subjects — stricter and more clinically honest evaluation",
    "PA2 vs PA3 (~1°C) is documented as the hardest adjacent pair; papers reporting 87% benefit from non-reactive subject majority",
    "CrossMod-Transformer uses full transformer stacks — parameters >> appropriate for N=67",
    "Our LOSO tests true cross-subject generalization; some published evaluations use random splits that leak subject identity",
    "65.3% on a rigorous LOSO-67 reactive-only benchmark is defensible and reproducible",
]
bullet_block(sl, defenses,
             Inches(0.5), y0 + Inches(3.48), W - Inches(0.9), Inches(3.0), size=14)
footer(sl, 17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Clinical Significance & Deployment
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Clinical Significance & Deployment Readiness")

# SHAP → Hardware implications
add_rect(sl, Inches(0.35), y0 + Inches(0.2), Inches(6.0), Inches(5.45), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "SHAP → HARDWARE REQUIREMENTS", Inches(0.5), y0 + Inches(0.35), Inches(5.5))
hw = [
    ("Priority 1", "GSR sensor", "Wearable sweat sensor / ICU adhesive patch — #1 feature by SHAP margin"),
    ("Priority 2", "ECG / PPG",  "Existing ICU cardiac monitor or optical pulse sensor — contributes ecg_max, ecg_shannon"),
    ("Priority 3", "Camera",     "Bedside 25fps camera — captures mouth & brow variability (top landmark features)"),
    ("Priority 4", "EMG",        "Electrode placement on trapezius — useful but logistically harder"),
]
hy = y0 + Inches(0.72)
for pri, hw_name, desc in hw:
    add_rect(sl, Inches(0.5), hy, Inches(1.0), Inches(0.52), fill=TEAL)
    add_text(sl, pri, Inches(0.55), hy + Inches(0.08), Inches(0.9), Inches(0.4), bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, hw_name, Inches(1.58), hy + Inches(0.08), Inches(1.5), Inches(0.4), bold=True, size=14, color=NAVY)
    add_text(sl, desc, Inches(3.15), hy + Inches(0.1), Inches(3.2), Inches(0.38), size=12, color=BLACK)
    hy += Inches(0.62)

section_label(sl, "DEPLOYMENT PROPERTIES", Inches(0.5), y0 + Inches(3.4), Inches(5))
props = ["✅  Non-invasive — no needles, no physical intervention",
         "✅  Real-time capable — feature extraction < 1 second per 5.5s window",
         "✅  Explainable via SHAP — clinician sees exactly why each prediction was made",
         "✅  Multimodal — graceful degradation if one modality unavailable",
         "✅  Subject-independent — tested on unseen subjects via LOSO"]
bullet_block(sl, props, Inches(0.5), y0 + Inches(3.75), Inches(5.7), Inches(1.9), size=13, bullet="")

# Right — system pipeline
add_rect(sl, Inches(6.6), y0 + Inches(0.2), Inches(6.35), Inches(5.45), fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.75))
section_label(sl, "SYSTEM PIPELINE", Inches(6.75), y0 + Inches(0.35), Inches(5.5))
steps = [
    ("1", "Sense",    "GSR patch + ECG + camera acquire 5.5s window"),
    ("2", "Extract",  "35 biosignal features + 22 landmark features"),
    ("3", "Normalize","Person-specific normalization within window"),
    ("4", "Predict",  "RF_bio + RF_lm → LogReg meta-learner"),
    ("5", "Explain",  "SHAP values generated per prediction"),
    ("6", "Alert",    "PA3 detected → clinical alert + feature breakdown"),
]
py = y0 + Inches(0.72)
for num, step, desc in steps:
    add_rect(sl, Inches(6.75), py, Inches(0.45), Inches(0.65), fill=NAVY)
    add_text(sl, num, Inches(6.78), py + Inches(0.1), Inches(0.38), Inches(0.48), bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(7.25), py, Inches(1.05), Inches(0.65), fill=TEAL)
    add_text(sl, step, Inches(7.28), py + Inches(0.1), Inches(0.98), Inches(0.48), bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, Inches(8.38), py + Inches(0.12), Inches(4.35), Inches(0.48), size=13, color=BLACK)
    py += Inches(0.77)
footer(sl, 18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Tech Stack & HPC
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GREY_BG)
y0 = navy_header(sl, "Technical Stack & Implementation")

simple_table(sl,
    ["Component", "Tool / Library", "Purpose"],
    [["Biosignal features",   "neurokit2, numpy",               "HRV extraction, signal statistics"],
     ["Facial landmarks",     "MediaPipe FaceMesh",             "468-point facial landmark detection"],
     ["ML models",            "scikit-learn 1.7.2 (pinned)",    "RF, LogReg — baseline stacked fusion"],
     ["GNN",                  "PyTorch Geometric (GAT)",        "Graph attention on landmark graph"],
     ["DANN / CrossMod",      "PyTorch (autograd, MultiheadAttention)", "Adversarial domain adaptation + cross-modal attention"],
     ["Explainability",       "SHAP TreeExplainer",             "Out-of-fold feature attribution"],
     ["Visualization",        "matplotlib, seaborn",            "SHAP plots, confusion matrix, per-subject charts"],
     ["HPC",                  "Hofstra StarHPC — H100 80GB GPU","SLURM jobs, GNN training"],
     ["Demo app",             "Streamlit",                      "Interactive pain prediction demo"],
     ["Eval protocol",        "LOSO (LeaveOneGroupOut)",        "67 folds, true cross-subject generalization"]],
    Inches(0.35), y0 + Inches(0.2), W - Inches(0.7),
    font_size=13)

# Architecture decisions rationale
add_rect(sl, Inches(0.35), y0 + Inches(4.55), W - Inches(0.7), Inches(0.38), fill=TEAL)
add_text(sl, "KEY ENGINEERING DECISIONS", Inches(0.55), y0 + Inches(4.6), W - Inches(1.0), Inches(0.3), bold=True, size=14, color=WHITE)
bullet_block(sl, [
    "scikit-learn pinned to 1.7.2 — matches saved empath_model.pkl; upgrading without re-saving breaks production model",
    "GNN GPU struct cache: precomputed edge offsets on GPU — eliminated 33MB/batch transfer, reduced fold time from 5min → 90sec",
    "matplotlib.use('Agg') throughout — all plots run headless on SLURM compute nodes with no display",
    "Out-of-fold SHAP — every SHAP value from a test subject the model never trained on; zero leakage",
], Inches(0.5), y0 + Inches(4.98), W - Inches(0.9), Inches(1.55), size=13)
footer(sl, 19)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, 0, W, Inches(0.08), fill=TEAL)
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), fill=TEAL)

add_text(sl, "Conclusion", Inches(0.6), Inches(0.3), W - Inches(1.2), Inches(0.7),
         bold=True, size=36, color=WHITE, align=PP_ALIGN.LEFT)

# 3 achievement boxes
ach = [
    ("What Was Built",
     "Multimodal pain classifier achieving 65.3% LOSO-67 accuracy on BioVid PA2 vs PA3 — the hardest adjacent pain level pair — evaluated with true subject-independent cross-validation."),
    ("What Was Discovered",
     "GSR slope is the #1 pain biomarker. Facial variability > static geometry. RF is the empirical ceiling at N=67. DANN reduces variance but not mean accuracy. PA3 is a held expression, not a fast movement."),
    ("Novel Contributions",
     "GNN on MediaPipe FaceMesh (first in BioVid literature) • DANN adversarial domain adaptation • CrossMod bidirectional attention • Landmark velocity features • Out-of-fold SHAP XAI layer"),
]
bw3 = Inches(4.0)
for i, (t, b) in enumerate(ach):
    x = Inches(0.4) + i * (bw3 + Inches(0.25))
    add_rect(sl, x, Inches(1.15), bw3, Inches(2.55), fill=TEAL)
    add_text(sl, t, x + Inches(0.15), Inches(1.2), bw3 - Inches(0.25), Inches(0.5),
             bold=True, size=15, color=WHITE)
    add_text(sl, b, x + Inches(0.15), Inches(1.75), bw3 - Inches(0.25), Inches(1.85),
             size=13, color=WHITE)

add_rect(sl, Inches(0.4), Inches(3.9), W - Inches(0.8), Inches(1.35), fill=WHITE)
add_text(sl, '"I built a subject-independent pain detection system using biosignals and facial '
             'landmarks, implemented and ablated GNNs, adversarial domain adaptation, and cross-modal '
             'attention — found that Random Forest remains the empirical ceiling at 67 subjects — '
             'and added SHAP explainability so clinicians can trust the predictions."',
         Inches(0.6), Inches(3.98), W - Inches(1.2), Inches(1.18),
         size=14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# Final metrics row
metrics = [
    ("65.3%",   "LOSO-67\nAccuracy"),
    ("15+",     "Architectures\nEvaluated"),
    ("5",       "Novel\nContributions"),
    ("67",      "Subjects\nTested"),
    ("2,680",   "Total\nSamples"),
    ("7",       "XAI\nPlots Generated"),
]
mx = Inches(0.4); my = Inches(5.42); mw = Inches(2.05); mh = Inches(1.78)
for val, lbl in metrics:
    add_rect(sl, mx, my, mw, mh, fill=RGBColor(0x0A, 0x20, 0x45))
    add_text(sl, val, mx + Inches(0.05), my + Inches(0.18), mw - Inches(0.08), Inches(0.88),
             bold=True, size=32, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, mx + Inches(0.05), my + Inches(1.05), mw - Inches(0.08), Inches(0.62),
             size=12, color=RGBColor(0xAA,0xCC,0xDD), align=PP_ALIGN.CENTER)
    mx += mw + Inches(0.12)
footer(sl, 20)


# ── Save ───────────────────────────────────────────────────────────────────────
OUT = "/Users/komalabelursrinivas/Desktop/Capstone/EmPath_v2/EmPath_v2_Presentation.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
